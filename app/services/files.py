"""Helpers for extracting text from uploaded documents."""
from pathlib import Path


def extract_text(path: str) -> str:
    """Best-effort text extraction from PDF / DOCX / PPTX / TXT."""
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(str(p))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        if ext == ".docx":
            from docx import Document
            doc = Document(str(p))
            return "\n".join(par.text for par in doc.paragraphs)
        if ext in {".pptx", ".ppt"}:
            from pptx import Presentation
            prs = Presentation(str(p))
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        chunks.append(shape.text)
            return "\n".join(chunks)
        if ext in {".txt", ".md"}:
            return p.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[Extraction error: {e}]"
    return "[Unsupported file type]"
